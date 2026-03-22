"""
SERP Evolution Analyzer — Five Blocks Internal Tool
Upload an IMPACT CSV, set a pivot date, get a branded PPTX deck + 2 Excel files.
"""

import streamlit as st
import pandas as pd
import re
import io
import os
import math
import base64
import tempfile
from datetime import datetime, date
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

# ─────────────────────────────────────────────
# CONSTANTS
# ─────────────────────────────────────────────

NAVY = RGBColor(0x1A, 0x36, 0x5D)
TEAL = RGBColor(0x0D, 0x7C, 0x7D)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
DARK = RGBColor(0x1F, 0x29, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xFA, 0xFA, 0xFA)

NAVY_HEX = "1A365D"
TEAL_HEX = "0D7C7D"
GRAY_HEX = "9CA3AF"
DARK_HEX = "1F2937"

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

DOMAIN_CATEGORIES = {
    'bloomberg.com': 'News', 'wsj.com': 'News', 'financialpost.com': 'News',
    'finance.yahoo.com': 'News', 'morningstar.com': 'News', 'costar.com': 'News',
    'theglobeandmail.com': 'News', 'bisnow.com': 'News', 'pionline.com': 'News',
    'internationalfinance.com': 'News', 'stocktitan.net': 'News',
    'institutionalinvestor.com': 'News', 'fintool.com': 'News', 'reuters.com': 'News',
    'cnbc.com': 'News', 'ft.com': 'News', 'nytimes.com': 'News', 'forbes.com': 'News',
    'businessinsider.com': 'News', 'barrons.com': 'News',
    'linkedin.com': 'Social', 'uk.linkedin.com': 'Social', 'instagram.com': 'Social',
    'x.com': 'Social', 'twitter.com': 'Social', 'facebook.com': 'Social',
    'youtube.com': 'Video', 'dailymotion.com': 'Video', 'vimeo.com': 'Video',
    'omny.fm': 'Podcast', 'fs.blog': 'Podcast', 'podcasts.apple.com': 'Podcast',
    'open.spotify.com': 'Podcast',
    'wallstreetoasis.com': 'Forum / UGC', 'reddit.com': 'Forum / UGC',
    'quora.com': 'Forum / UGC',
    'marketscreener.com': 'Financial Data',
    'sec.gov': 'Regulatory',
    'find-and-update.company-information.service.gov.uk': 'Regulatory',
}

VALID_CATEGORIES = [
    'News', 'Corporate / Owned', 'Social', 'Video', 'Podcast',
    'Forum / UGC', 'Research', 'Financial Data', 'Conference / Event',
    'Regulatory', 'Other'
]

# ─────────────────────────────────────────────
# LOGO LOADING
# ─────────────────────────────────────────────

@st.cache_resource
def _get_logo_paths():
    """
    Load logos from Streamlit secrets (base64-encoded) or fall back to local assets/ folder.
    Returns (logo_path, logo_white_path) as temp file paths.
    """
    logo_path = None
    logo_white_path = None

    # Try secrets first
    try:
        logo_b64 = st.secrets.get("LOGO_PNG", None)
        logo_white_b64 = st.secrets.get("LOGO_WHITE_PNG", None)

        if logo_b64:
            tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            tmp.write(base64.b64decode(logo_b64))
            tmp.close()
            logo_path = tmp.name

        if logo_white_b64:
            tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            tmp.write(base64.b64decode(logo_white_b64))
            tmp.close()
            logo_white_path = tmp.name
    except Exception:
        pass

    # Fall back to local assets/ folder (for local dev)
    if not logo_path:
        local = os.path.join(os.path.dirname(__file__), 'assets', 'logo.png')
        if os.path.exists(local):
            logo_path = local

    if not logo_white_path:
        local = os.path.join(os.path.dirname(__file__), 'assets', 'logo_white.png')
        if os.path.exists(local):
            logo_white_path = local

    return logo_path, logo_white_path

# ─────────────────────────────────────────────
# STEP 1: DATA PROCESSING
# ─────────────────────────────────────────────

def extract_domain(url):
    m = re.search(r'https?://(?:www\.)?([^/]+)', str(url))
    return m.group(1) if m else str(url)


def make_url_label(url, all_urls_for_domain=None):
    """Create a short display label for a URL, disambiguating when needed."""
    domain = extract_domain(url)
    if not all_urls_for_domain or len(all_urls_for_domain) <= 1:
        return domain

    path = re.sub(r'https?://[^/]+', '', str(url)).strip('/')
    if not path:
        return domain

    # Clean up the path into a readable hint
    # Remove common prefixes and file extensions
    path = re.sub(r'\.(html?|php|aspx?|jsp|pdf)$', '', path, flags=re.I)
    parts = [p for p in path.split('/') if p]
    if not parts:
        return domain

    # Try to extract a meaningful short hint
    last = parts[-1]
    # URL-decode
    last = last.replace('%20', ' ').replace('%2F', '/').replace('%26', '&')
    # Clean query params for youtube etc
    if 'watch' in last and 'v=' in last:
        return f"{domain} (video {all_urls_for_domain.index(url) + 1})" if url in all_urls_for_domain else f"{domain} (video)"

    # For very long slugs, try to shorten
    last = re.sub(r'[-_]+', ' ', last)
    # If it's a clear content type indicator
    content_types = {
        'news': 'news', 'article': 'article', 'blog': 'blog', 'video': 'video',
        'podcast': 'podcast', 'interview': 'interview', 'bio': 'bio',
        'about': 'bio', 'leadership': 'bio', 'profile': 'profile',
        'press': 'press release', 'pdf': 'PDF',
    }
    path_lower = path.lower()
    for keyword, label in content_types.items():
        if keyword in path_lower:
            # Check if there are multiple of the same type for this domain
            same_type = [u for u in all_urls_for_domain if keyword in u.lower()]
            if len(same_type) > 1:
                idx = same_type.index(url) + 1 if url in same_type else 1
                return f"{domain} ({label} {idx})"
            return f"{domain} ({label})"

    # If it's a PDF file
    if '.pdf' in url.lower():
        return f"{domain} (PDF)"

    # Fallback: use a very short version of the last path segment
    if len(last) > 20:
        last = last[:18].rsplit(' ', 1)[0]
    return f"{domain} ({last.strip()})" if last.strip() else domain


def process_impact_csv(file_bytes, pivot_date):
    """Parse IMPACT CSV, filter to Standard results, split at pivot date."""
    try:
        df = pd.read_csv(io.BytesIO(file_bytes), encoding='utf-8')
    except UnicodeDecodeError:
        df = pd.read_csv(io.BytesIO(file_bytes), encoding='latin-1')

    df['parsed_date'] = pd.to_datetime(df['Date'], format='%d %b %Y')
    std = df[df['Result Type'] == 'Standard'].copy()

    pivot_dt = pd.Timestamp(pivot_date)
    pre = std[std['parsed_date'] < pivot_dt]
    post = std[std['parsed_date'] >= pivot_dt]

    pre_days = sorted(pre['parsed_date'].dt.date.unique())
    post_days = sorted(post['parsed_date'].dt.date.unique())

    url_data = {}
    for url in std['URL'].unique():
        pre_dates = set(pre[pre['URL'] == url]['parsed_date'].dt.date.unique())
        post_dates = set(post[post['URL'] == url]['parsed_date'].dt.date.unique())
        pre_count = len(pre_dates)
        post_count = len(post_dates)
        if pre_count > 0 and post_count > 0:
            status = 'Persistent'
        elif post_count > 0:
            status = 'New'
        else:
            status = 'Dropped'
        domain = extract_domain(url)
        url_data[url] = {
            'url': url, 'domain': domain, 'status': status,
            'pre_days': pre_count, 'post_days': post_count,
        }

    return url_data, len(pre_days), len(post_days), pre, post


# ─────────────────────────────────────────────
# STEP 2: DOMAIN CATEGORIZATION
# ─────────────────────────────────────────────

def categorize_domains(url_data, owned_domains=None, api_key=None):
    """Categorize each URL's domain using built-in map + optional Claude API fallback."""
    owned = set()
    if owned_domains:
        for d in owned_domains.split(','):
            d = d.strip().lower()
            if d:
                owned.add(d)

    unmapped_urls = []
    for url, info in url_data.items():
        domain = info['domain'].lower()
        if domain in owned or any(domain.endswith('.' + o) for o in owned):
            info['category'] = 'Corporate / Owned'
        elif domain in DOMAIN_CATEGORIES:
            info['category'] = DOMAIN_CATEGORIES[domain]
        else:
            info['category'] = None
            unmapped_urls.append(url)

    if unmapped_urls and api_key:
        try:
            _classify_with_claude(url_data, unmapped_urls, api_key)
        except Exception as e:
            st.warning(f"Claude API classification failed: {e}. Unmapped domains set to 'Other'.")
            for url in unmapped_urls:
                if url_data[url]['category'] is None:
                    url_data[url]['category'] = 'Other'
    else:
        for url in unmapped_urls:
            if url_data[url]['category'] is None:
                url_data[url]['category'] = 'Other'

    return url_data


def _classify_with_claude(url_data, unmapped_urls, api_key):
    import anthropic
    client = anthropic.Anthropic(api_key=api_key)
    url_list = "\n".join(unmapped_urls)
    categories = ", ".join(VALID_CATEGORIES)
    msg = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=2000,
        temperature=0.0,
        messages=[{
            "role": "user",
            "content": f"""Classify each URL below into exactly one category.
Categories: {categories}

Consider the full URL path, not just the domain. For example:
- morningstar.com/news/ → News
- morningstar.com/people/ → Financial Data

Return ONLY lines in format: URL | Category

URLs:
{url_list}"""
        }]
    )
    text = msg.content[0].text
    for line in text.strip().split('\n'):
        if '|' in line:
            parts = line.split('|', 1)
            url_str = parts[0].strip()
            cat = parts[1].strip()
            if url_str in url_data and cat in VALID_CATEGORIES:
                url_data[url_str]['category'] = cat


# ─────────────────────────────────────────────
# STEP 3: SOURCE TYPE AGGREGATION
# ─────────────────────────────────────────────

def compute_source_types(url_data, pre_df, post_df, n_pre, n_post):
    """Compute slot-days per category per period."""
    url_to_cat = {u: d['category'] for u, d in url_data.items()}

    def count_slot_days(period_df):
        cats = defaultdict(int)
        domains = defaultdict(lambda: defaultdict(int))
        for _, row in period_df.iterrows():
            url = row['URL']
            cat = url_to_cat.get(url, 'Other')
            domain = extract_domain(url)
            cats[cat] += 1
            domains[cat][domain] += 1
        return cats, domains

    pre_cats, pre_domains = count_slot_days(pre_df)
    post_cats, post_domains = count_slot_days(post_df)

    all_cats = sorted(set(list(pre_cats.keys()) + list(post_cats.keys())),
                      key=lambda c: -(pre_cats.get(c, 0) + post_cats.get(c, 0)))

    cat_summary = []
    for cat in all_cats:
        pre_sd = pre_cats.get(cat, 0)
        post_sd = post_cats.get(cat, 0)
        cat_summary.append({
            'category': cat, 'pre_slot_days': pre_sd,
            'post_slot_days': post_sd, 'change': post_sd - pre_sd
        })

    all_domains_set = set()
    for d in list(pre_domains.values()) + list(post_domains.values()):
        all_domains_set.update(d.keys())

    domain_detail = []
    for cat in all_cats:
        cat_doms = set(list(pre_domains.get(cat, {}).keys()) + list(post_domains.get(cat, {}).keys()))
        for dom in sorted(cat_doms):
            pre_sd = pre_domains.get(cat, {}).get(dom, 0)
            post_sd = post_domains.get(cat, {}).get(dom, 0)
            domain_detail.append({
                'category': cat, 'domain': dom, 'pre_slot_days': pre_sd,
                'post_slot_days': post_sd, 'total': pre_sd + post_sd
            })

    return cat_summary, domain_detail


# ─────────────────────────────────────────────
# STEP 4: EXECUTIVE SUMMARY (Claude API)
# ─────────────────────────────────────────────

def generate_executive_summary(url_data, cat_summary, client_name, event_desc, n_pre, n_post, api_key):
    """Generate 4 executive summary points using Claude API."""
    new_urls = [(d['domain'], d['post_days']) for d in url_data.values() if d['status'] == 'New']
    new_urls.sort(key=lambda x: -x[1])
    persistent = [(d['domain'], d['pre_days'], d['post_days']) for d in url_data.values() if d['status'] == 'Persistent']
    persistent.sort(key=lambda x: -(x[1] + x[2]))
    dropped = [d['domain'] for d in url_data.values() if d['status'] == 'Dropped']

    n_new = len(new_urls)
    n_persistent = len(persistent)
    n_dropped = len(dropped)
    total = n_new + n_persistent + n_dropped

    cat_shifts = "\n".join(
        f"  {c['category']}: {c['pre_slot_days']} → {c['post_slot_days']} (change: {c['change']:+d})"
        for c in cat_summary
    )
    top_new = "\n".join(f"  {d}: {days} days" for d, days in new_urls[:5])
    top_persistent = "\n".join(f"  {d}: {pre}→{post}" for d, pre, post in persistent[:5])

    import anthropic
    client = anthropic.Anthropic(api_key=api_key)
    msg = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=1000,
        temperature=0.3,
        messages=[{
            "role": "user",
            "content": f"""Generate exactly 4 executive summary points about SERP changes for {client_name} around: {event_desc}.

Data:
- Pre period: {n_pre} days, Post period: {n_post} days
- {total} unique URLs: {n_new} new, {n_persistent} persistent, {n_dropped} dropped

Top new URLs (by days present in post):
{top_new}

Top persistent URLs (pre→post days):
{top_persistent}

Dropped URLs: {', '.join(dropped) if dropped else 'None significant'}

Slot-day shifts by source category:
{cat_shifts}

Format each point as:
HEADLINE: [bold one-sentence headline]
DETAIL: [1-2 sentences with specific numbers]

Be concise, data-driven, and insightful. Focus on what matters for reputation management. Do not recommend Knowledge Panel claiming/verification."""
        }]
    )

    text = msg.content[0].text
    points = _parse_summary_points(text)

    # If Claude's response couldn't be parsed, fall back to data-driven templates
    if len(points) < 2:
        return generate_fallback_summary(url_data, cat_summary, client_name, event_desc, n_pre, n_post)

    # Pad with fallback points if we got some but fewer than 4
    if len(points) < 4:
        fallback = generate_fallback_summary(url_data, cat_summary, client_name, event_desc, n_pre, n_post)
        for fb in fallback:
            if len(points) >= 4:
                break
            if fb[0] not in [p[0] for p in points]:
                points.append(fb)

    return points[:4]


def _parse_summary_points(text):
    """Parse executive summary points from Claude's response, handling varied formats."""
    points = []

    # Strategy 1: Explicit HEADLINE: / DETAIL: format
    current_headline = None
    current_detail = None
    for line in text.strip().split('\n'):
        line = line.strip()
        if not line:
            continue
        if line.upper().startswith('HEADLINE:'):
            if current_headline:
                points.append((current_headline, current_detail or ''))
            current_headline = line.split(':', 1)[1].strip().strip('*')
            current_detail = None
        elif line.upper().startswith('DETAIL:'):
            current_detail = line.split(':', 1)[1].strip()
    if current_headline:
        points.append((current_headline, current_detail or ''))

    if len(points) >= 2:
        return points

    # Strategy 2: Numbered points with bold headline — e.g. "1. **Headline here**\nDetail text"
    # Also handles "1. Headline here\nDetail text" without bold
    points = []
    lines = text.strip().split('\n')
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        # Match numbered line: "1.", "1)", "- ", "• "
        m = re.match(r'^(?:\d+[\.\)]\s*|[-•]\s*)(.*)', line)
        if m:
            headline = m.group(1).strip()
            # Strip markdown bold
            headline = re.sub(r'\*\*(.+?)\*\*', r'\1', headline)
            headline = headline.strip('*').strip()
            # Collect subsequent non-numbered lines as detail
            detail_lines = []
            i += 1
            while i < len(lines):
                next_line = lines[i].strip()
                if not next_line:
                    i += 1
                    continue
                if re.match(r'^(?:\d+[\.\)]\s*|[-•]\s*)', next_line):
                    break
                detail_lines.append(next_line)
                i += 1
            detail = ' '.join(detail_lines)
            detail = re.sub(r'\*\*(.+?)\*\*', r'\1', detail).strip()
            if headline:
                points.append((headline, detail))
        else:
            i += 1

    if len(points) >= 2:
        return points

    # Strategy 3: Bold lines as headlines, following lines as detail
    points = []
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        bold_match = re.match(r'^\*\*(.+?)\*\*[:\s]*(.*)', line)
        if bold_match:
            headline = bold_match.group(1).strip()
            rest = bold_match.group(2).strip()
            detail_lines = [rest] if rest else []
            i += 1
            while i < len(lines):
                next_line = lines[i].strip()
                if not next_line:
                    i += 1
                    continue
                if re.match(r'^\*\*', next_line):
                    break
                detail_lines.append(next_line)
                i += 1
            detail = ' '.join(detail_lines).strip()
            if headline:
                points.append((headline, detail))
        else:
            i += 1

    return points


def generate_fallback_summary(url_data, cat_summary, client_name, event_desc, n_pre, n_post):
    """Generate basic summary points without Claude API."""
    new_urls = sorted(
        [(d['domain'], d['post_days']) for d in url_data.values() if d['status'] == 'New'],
        key=lambda x: -x[1]
    )
    n_new = sum(1 for d in url_data.values() if d['status'] == 'New')
    n_persistent = sum(1 for d in url_data.values() if d['status'] == 'Persistent')
    n_dropped = sum(1 for d in url_data.values() if d['status'] == 'Dropped')
    total = n_new + n_persistent + n_dropped

    top_new_str = ", ".join(d for d, _ in new_urls[:3]) if new_urls else "none"

    biggest_gain = max(cat_summary, key=lambda c: c['change']) if cat_summary else None
    biggest_loss = min(cat_summary, key=lambda c: c['change']) if cat_summary else None

    points = [
        (f"{n_new} new URLs entered the SERP after the event.",
         f"Top new entrants: {top_new_str}. These appeared only in the post-announcement period."),
        (f"{n_persistent} URLs persisted across both periods.",
         f"Out of {total} total unique URLs, {n_persistent} appeared in both the pre and post periods."),
    ]
    if biggest_gain:
        points.append((
            f"{biggest_gain['category']} saw the largest gain.",
            f"Slot-days went from {biggest_gain['pre_slot_days']} to {biggest_gain['post_slot_days']} ({biggest_gain['change']:+d} change)."
        ))
    if biggest_loss and biggest_loss['change'] < 0:
        points.append((
            f"{biggest_loss['category']} was most displaced.",
            f"Slot-days dropped from {biggest_loss['pre_slot_days']} to {biggest_loss['post_slot_days']} ({biggest_loss['change']:+d} change)."
        ))

    while len(points) < 4:
        points.append(("Dropped URLs were minor.", f"{n_dropped} dropped URLs were typically single-day appearances."))
    return points[:4]


# ─────────────────────────────────────────────
# STEP 5: EXCEL GENERATION
# ─────────────────────────────────────────────

def generate_url_analysis_excel(url_data, n_pre, n_post, filename, pivot_date):
    wb = Workbook()
    ws = wb.active
    ws.title = "URL Analysis"

    header_font = Font(name='Arial', bold=True, color='FFFFFF', size=11)
    header_fill = PatternFill('solid', fgColor=NAVY_HEX)
    new_fill = PatternFill('solid', fgColor='E6F4EA')
    dropped_fill = PatternFill('solid', fgColor='FDE8E8')

    headers = ['URL', 'Domain', 'Status', f'Pre: Days Present (of {n_pre})', f'Post: Days Present (of {n_post})']
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center')

    sorted_urls = sorted(url_data.values(), key=lambda d: (
        {'Persistent': 0, 'New': 1, 'Dropped': 2}[d['status']],
        -(d['pre_days'] + d['post_days'])
    ))

    for i, info in enumerate(sorted_urls, 2):
        ws.cell(row=i, column=1, value=info['url'])
        ws.cell(row=i, column=2, value=info['domain'])
        ws.cell(row=i, column=3, value=info['status'])
        ws.cell(row=i, column=4, value=info['pre_days'])
        ws.cell(row=i, column=5, value=info['post_days'])
        if info['status'] == 'New':
            for c in range(1, 6):
                ws.cell(row=i, column=c).fill = new_fill
        elif info['status'] == 'Dropped':
            for c in range(1, 6):
                ws.cell(row=i, column=c).fill = dropped_fill

    ws.column_dimensions['A'].width = 80
    ws.column_dimensions['B'].width = 30
    ws.column_dimensions['C'].width = 14
    ws.column_dimensions['D'].width = 24
    ws.column_dimensions['E'].width = 24
    ws.freeze_panes = 'A2'

    # Summary sheet
    ws2 = wb.create_sheet("Summary")
    n_new = sum(1 for d in url_data.values() if d['status'] == 'New')
    n_persistent = sum(1 for d in url_data.values() if d['status'] == 'Persistent')
    n_dropped = sum(1 for d in url_data.values() if d['status'] == 'Dropped')
    total = len(url_data)

    summary_data = [
        ('Input file', filename),
        ('Event date', str(pivot_date)),
        ('Pre-announcement days', n_pre),
        ('Post-announcement days', n_post),
        ('Total unique URLs', total),
        ('Persistent', n_persistent),
        ('New', n_new),
        ('Dropped', n_dropped),
    ]
    for i, (metric, val) in enumerate(summary_data, 1):
        ws2.cell(row=i, column=1, value=metric).font = Font(name='Arial', bold=True)
        ws2.cell(row=i, column=2, value=val)
    ws2.column_dimensions['A'].width = 28
    ws2.column_dimensions['B'].width = 50

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def generate_source_type_excel(cat_summary, domain_detail, n_pre, n_post):
    wb = Workbook()
    ws = wb.active
    ws.title = "Category Summary"

    header_font = Font(name='Arial', bold=True, color='FFFFFF', size=11)
    header_fill = PatternFill('solid', fgColor=NAVY_HEX)

    headers = ['Category', f'Pre Slot-Days ({n_pre}d)', f'Post Slot-Days ({n_post}d)', 'Change']
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center')

    for i, row in enumerate(cat_summary, 2):
        ws.cell(row=i, column=1, value=row['category'])
        ws.cell(row=i, column=2, value=row['pre_slot_days'])
        ws.cell(row=i, column=3, value=row['post_slot_days'])
        change_cell = ws.cell(row=i, column=4, value=row['change'])
        if row['change'] > 0:
            change_cell.font = Font(name='Arial', color='0D7C7D')
        elif row['change'] < 0:
            change_cell.font = Font(name='Arial', color='CC0000')

    ws.column_dimensions['A'].width = 24
    ws.column_dimensions['B'].width = 22
    ws.column_dimensions['C'].width = 22
    ws.column_dimensions['D'].width = 12

    # Domain Detail sheet
    ws2 = wb.create_sheet("Domain Detail")
    headers2 = ['Category', 'Domain', f'Pre Slot-Days ({n_pre}d)', f'Post Slot-Days ({n_post}d)', 'Total']
    for col, h in enumerate(headers2, 1):
        cell = ws2.cell(row=1, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center')

    for i, row in enumerate(domain_detail, 2):
        ws2.cell(row=i, column=1, value=row['category'])
        ws2.cell(row=i, column=2, value=row['domain'])
        ws2.cell(row=i, column=3, value=row['pre_slot_days'])
        ws2.cell(row=i, column=4, value=row['post_slot_days'])
        ws2.cell(row=i, column=5, value=row['total'])

    ws2.column_dimensions['A'].width = 24
    ws2.column_dimensions['B'].width = 36
    ws2.column_dimensions['C'].width = 22
    ws2.column_dimensions['D'].width = 22
    ws2.column_dimensions['E'].width = 12

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────
# STEP 6: PPTX GENERATION (python-pptx)
# ─────────────────────────────────────────────

def _format_date(d, fmt='%b %-d, %Y'):
    """Format a date, handling platform differences for %-d."""
    try:
        return d.strftime(fmt)
    except ValueError:
        return d.strftime(fmt.replace('%-d', '%d').replace('%-m', '%m'))


def _add_textbox(slide, left, top, width, height, text, font_name='Calibri',
                 font_size=12, color=DARK, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    # Set margins to minimal
    txBox.text_frame.margin_left = Emu(0)
    txBox.text_frame.margin_right = Emu(0)
    txBox.text_frame.margin_top = Emu(0)
    txBox.text_frame.margin_bottom = Emu(0)
    return txBox


def _add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_line(slide, left, top, width, color, thickness=Pt(2)):
    shape = slide.shapes.add_shape(1, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_logo(slide, logo_path, is_dark=False):
    """Add Five Blocks logo bottom-right."""
    try:
        img = slide.shapes.add_picture(
            logo_path, Inches(8.3), Inches(4.85), Inches(1.4), Inches(0.5)
        )
    except Exception:
        pass


def _add_hyperlink(run, url):
    """Add hyperlink to a run."""
    rPr = run._r.get_or_add_rPr()
    hlinkClick = rPr.makeelement(qn('a:hlinkClick'), {})
    hlinkClick.set(qn('r:id'), '')
    rPr.append(hlinkClick)
    # Use the relationship approach
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        rel = run.part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
        hlinkClick.set(qn('r:id'), rel)
    except Exception:
        pass


def build_slide_1(prs, client_name, event_desc, date_range_str, logo_white_path):
    """Title slide — navy background."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY

    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(8), Inches(1.2),
                 client_name, 'Calibri Light', 42, WHITE, bold=False)

    _add_textbox(slide, Inches(0.8), Inches(2.5), Inches(8), Inches(0.5),
                 "Search result evolution", 'Calibri Light', 18, GRAY)

    _add_textbox(slide, Inches(0.8), Inches(3.2), Inches(8), Inches(0.8),
                 f"{event_desc}\n{date_range_str}", 'Calibri', 13, GRAY)

    _add_logo(slide, logo_white_path, is_dark=True)


def build_slide_2(prs, summary_points, logo_path):
    """Executive summary — white background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 "Executive summary", 'Calibri Light', 26, NAVY)

    # Teal accent line under title
    _add_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), TEAL, Pt(3))

    y_start = Inches(1.3)
    for i, (headline, detail) in enumerate(summary_points):
        y = y_start + Inches(i * 1.0)
        num_str = f"0{i+1}"
        _add_textbox(slide, Inches(0.6), y, Inches(0.6), Inches(0.5),
                     num_str, 'Calibri Light', 24, RGBColor(0xD1, 0xD5, 0xDB))

        _add_textbox(slide, Inches(1.2), y, Inches(8.2), Inches(0.35),
                     headline, 'Calibri', 13, DARK, bold=True)

        _add_textbox(slide, Inches(1.2), y + Inches(0.35), Inches(8.2), Inches(0.55),
                     detail, 'Calibri', 11, GRAY)

    _add_logo(slide, logo_path)


def build_slide_3(prs, url_data, n_pre, n_post, client_name, event_desc, pivot_date, logo_path):
    """Stats cards — light gray background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG

    n_new = sum(1 for d in url_data.values() if d['status'] == 'New')
    n_persistent = sum(1 for d in url_data.values() if d['status'] == 'Persistent')
    n_dropped = sum(1 for d in url_data.values() if d['status'] == 'Dropped')
    total = n_new + n_persistent + n_dropped

    _add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8.5), Inches(0.7),
                 "The announcement reshaped search", 'Calibri Light', 26, DARK)

    pivot_str = _format_date(pivot_date, '%B %-d') if hasattr(pivot_date, 'strftime') else str(pivot_date)
    _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(8.5), Inches(0.6),
                 f"Of {total} unique URLs appearing as standard results across the full period, the majority entered the SERP only after the {event_desc} on {pivot_str}.",
                 'Calibri', 12, GRAY)

    cards = [
        (n_new, "New URLs", "Appeared only after the\nannouncement", TEAL),
        (n_persistent, "Persistent URLs", "Present in both pre and post\nperiods", NAVY),
        (n_dropped, "Dropped URLs", "Disappeared after the\nannouncement", GRAY),
    ]

    card_w = Inches(2.5)
    gap = Inches(0.4)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2

    for idx, (num, label, desc, accent_color) in enumerate(cards):
        x = start_x + (card_w + gap) * idx
        y = Inches(1.8)

        # Accent line
        _add_line(slide, x + Inches(0.2), y, card_w - Inches(0.4), accent_color, Pt(3))

        # Number
        _add_textbox(slide, x, y + Inches(0.2), card_w, Inches(0.8),
                     str(num), 'Calibri Light', 48, accent_color, alignment=PP_ALIGN.CENTER)

        # Label
        _add_textbox(slide, x, y + Inches(1.0), card_w, Inches(0.35),
                     label, 'Calibri', 14, DARK, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        _add_textbox(slide, x, y + Inches(1.35), card_w, Inches(0.6),
                     desc, 'Calibri', 11, GRAY, alignment=PP_ALIGN.CENTER)

    # Footnote about dropped
    if n_dropped > 0 and all(url_data[u]['pre_days'] <= 2 for u in url_data if url_data[u]['status'] == 'Dropped'):
        _add_textbox(slide, Inches(0.8), Inches(4.1), Inches(8), Inches(0.3),
                     f"{n_dropped} dropped URL{'s' if n_dropped != 1 else ''} {'were' if n_dropped != 1 else 'was'} all single-day appearances — not significant losses.",
                     'Calibri', 11, GRAY, italic=True)

    _add_logo(slide, logo_path)


def build_slide_4_v2(prs, url_data, n_pre, n_post, logo_path):
    """Persistent URLs — butterfly chart (revised layout matching reference)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    persistent = [d for d in url_data.values() if d['status'] == 'Persistent']
    persistent.sort(key=lambda d: -(d['pre_days'] + d['post_days']))

    _add_textbox(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.6),
                 "Persistent URLs: pre vs. post", 'Calibri Light', 26, DARK)
    _add_textbox(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.35),
                 f"Days present out of {n_post} in each period", 'Calibri', 12, GRAY)

    # Legend
    _add_rect(slide, Inches(0.8), Inches(1.18), Inches(0.2), Inches(0.14), NAVY)
    _add_textbox(slide, Inches(1.05), Inches(1.14), Inches(0.5), Inches(0.25), "Pre", 'Calibri', 10, GRAY)
    _add_rect(slide, Inches(1.5), Inches(1.18), Inches(0.2), Inches(0.14), TEAL)
    _add_textbox(slide, Inches(1.75), Inches(1.14), Inches(0.5), Inches(0.25), "Post", 'Calibri', 10, GRAY)

    if not persistent:
        _add_textbox(slide, Inches(2), Inches(2.5), Inches(6), Inches(0.5),
                     "No persistent URLs found.", 'Calibri', 14, GRAY, alignment=PP_ALIGN.CENTER)
        _add_logo(slide, logo_path)
        return

    max_days = max(max(d['pre_days'], d['post_days']) for d in persistent)
    n_rows = len(persistent)

    chart_top = Inches(1.45)
    available_h = Inches(3.6)
    row_h = min(Inches(0.28), available_h / n_rows)
    bar_h = row_h * 0.6

    # Layout zones (matching reference: pre bars left, labels center, post bars right)
    pre_bar_right = Inches(4.1)  # right edge of pre bars
    max_pre_bar_w = Inches(3.0)
    label_center = Inches(5.0)  # center of labels
    label_w = Inches(1.7)
    post_bar_left = Inches(5.85)  # left edge of post bars
    max_post_bar_w = Inches(3.3)

    # Build disambiguation
    domain_urls = defaultdict(list)
    for d in persistent:
        domain_urls[d['domain']].append(d['url'])

    for i, d in enumerate(persistent):
        y = chart_top + row_h * i
        bar_y = y + (row_h - bar_h) / 2
        label = make_url_label(d['url'], domain_urls.get(d['domain']))

        # Pre bar (right-aligned to pre_bar_right)
        pre_w = max(Emu(1), int(max_pre_bar_w * d['pre_days'] / max_days))
        pre_left = pre_bar_right - pre_w
        _add_rect(slide, pre_left, bar_y, pre_w, bar_h, NAVY)

        # Pre value (left of bar)
        _add_textbox(slide, pre_left - Inches(0.35), y, Inches(0.32), row_h,
                     str(d['pre_days']), 'Calibri', 8, GRAY, alignment=PP_ALIGN.RIGHT)

        # Center label (hyperlinked to URL)
        txb = _add_textbox(slide, label_center - label_w / 2, y, label_w, row_h,
                           label, 'Calibri', 8, GRAY, alignment=PP_ALIGN.CENTER)
        try:
            _add_hyperlink(txb.text_frame.paragraphs[0].runs[0], d['url'])
        except Exception:
            pass

        # Post bar
        post_w = max(Emu(1), int(max_post_bar_w * d['post_days'] / max_days))
        _add_rect(slide, post_bar_left, bar_y, post_w, bar_h, TEAL)

        # Post value (right of bar)
        _add_textbox(slide, post_bar_left + post_w + Emu(20000), y, Inches(0.35), row_h,
                     str(d['post_days']), 'Calibri', 8, GRAY)

    _add_logo(slide, logo_path)


def build_slide_5(prs, url_data, n_post, logo_path):
    """New URLs — horizontal bar chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    new_urls = [d for d in url_data.values() if d['status'] == 'New']
    new_urls.sort(key=lambda d: -d['post_days'])
    total_new = len(new_urls)
    show = new_urls[:10]
    remaining = total_new - len(show)

    _add_textbox(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.6),
                 "New URLs: post-announcement", 'Calibri Light', 26, DARK)
    _add_textbox(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.35),
                 f"{total_new} URLs appeared only after the event. Top {min(10, total_new)} by days present (of {n_post}).",
                 'Calibri', 12, GRAY)

    if not show:
        _add_textbox(slide, Inches(2), Inches(2.5), Inches(6), Inches(0.5),
                     "No new URLs found.", 'Calibri', 14, GRAY, alignment=PP_ALIGN.CENTER)
        _add_logo(slide, logo_path)
        return

    max_days = show[0]['post_days'] if show else 1
    n_rows = len(show)
    chart_top = Inches(1.3)
    available_h = Inches(3.5)
    row_h = min(Inches(0.32), available_h / n_rows)
    bar_h = row_h * 0.55

    label_right = Inches(3.2)
    bar_left = Inches(3.4)
    max_bar_w = Inches(5.5)

    domain_urls = defaultdict(list)
    for d in new_urls:
        domain_urls[d['domain']].append(d['url'])

    for i, d in enumerate(show):
        y = chart_top + row_h * i
        bar_y = y + (row_h - bar_h) / 2
        label = make_url_label(d['url'], domain_urls.get(d['domain']))

        # Label (right-aligned, hyperlinked to URL)
        txb = _add_textbox(slide, Inches(0.3), y, label_right - Inches(0.4), row_h,
                     label, 'Calibri', 10, GRAY, alignment=PP_ALIGN.RIGHT)
        try:
            _add_hyperlink(txb.text_frame.paragraphs[0].runs[0], d['url'])
        except Exception:
            pass

        # Bar
        bar_w = max(Emu(1), int(max_bar_w * d['post_days'] / max_days))
        _add_rect(slide, bar_left, bar_y, bar_w, bar_h, TEAL)

        # Value
        _add_textbox(slide, bar_left + bar_w + Emu(30000), y, Inches(0.4), row_h,
                     str(d['post_days']), 'Calibri', 9, GRAY)

    # Footnote
    if remaining > 0:
        foot_y = chart_top + row_h * len(show) + Inches(0.15)
        _add_textbox(slide, Inches(0.8), foot_y, Inches(8), Inches(0.3),
                     f"+ {remaining} additional URLs appeared 1–2 times (transient news cycle coverage).",
                     'Calibri', 11, GRAY, italic=True)

    _add_logo(slide, logo_path)


def build_slide_6(prs, cat_summary, n_pre, n_post, logo_path):
    """Source type shift — grouped horizontal bar chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(slide, Inches(0.8), Inches(0.25), Inches(8), Inches(0.6),
                 "Source type shift", 'Calibri Light', 26, DARK)
    _add_textbox(slide, Inches(0.8), Inches(0.8), Inches(8), Inches(0.35),
                 f"Slot-days by source category across {n_pre}-day pre- and post-announcement periods.",
                 'Calibri', 12, GRAY)

    # Legend
    _add_rect(slide, Inches(0.8), Inches(1.18), Inches(0.2), Inches(0.14), NAVY)
    _add_textbox(slide, Inches(1.05), Inches(1.14), Inches(0.5), Inches(0.25), "Pre", 'Calibri', 10, GRAY)
    _add_rect(slide, Inches(1.5), Inches(1.18), Inches(0.2), Inches(0.14), TEAL)
    _add_textbox(slide, Inches(1.75), Inches(1.14), Inches(0.5), Inches(0.25), "Post", 'Calibri', 10, GRAY)

    # Sort by total descending
    sorted_cats = sorted(cat_summary, key=lambda c: -(c['pre_slot_days'] + c['post_slot_days']))

    if not sorted_cats:
        _add_logo(slide, logo_path)
        return

    max_sd = max(max(c['pre_slot_days'], c['post_slot_days']) for c in sorted_cats)
    n_cats = len(sorted_cats)
    chart_top = Inches(1.45)
    available_h = Inches(3.3)
    group_h = min(Inches(0.5), available_h / n_cats)
    bar_h = group_h * 0.35
    gap = group_h * 0.08

    label_right = Inches(2.2)
    bar_left = Inches(2.4)
    max_bar_w = Inches(6.5)

    for i, c in enumerate(sorted_cats):
        y = chart_top + group_h * i

        # Category label
        _add_textbox(slide, Inches(0.3), y, label_right - Inches(0.4), group_h,
                     c['category'], 'Calibri', 10, DARK, alignment=PP_ALIGN.RIGHT)

        # Pre bar (top)
        pre_bar_y = y + gap
        pre_w = max(Emu(1), int(max_bar_w * c['pre_slot_days'] / max_sd)) if c['pre_slot_days'] > 0 else Emu(1)
        _add_rect(slide, bar_left, pre_bar_y, pre_w, bar_h, NAVY)
        if c['pre_slot_days'] > 0:
            _add_textbox(slide, bar_left + pre_w + Emu(30000), pre_bar_y - Emu(10000), Inches(0.4), bar_h + Emu(20000),
                         str(c['pre_slot_days']), 'Calibri', 8, GRAY)

        # Post bar (bottom)
        post_bar_y = pre_bar_y + bar_h + gap
        post_w = max(Emu(1), int(max_bar_w * c['post_slot_days'] / max_sd)) if c['post_slot_days'] > 0 else Emu(1)
        _add_rect(slide, bar_left, post_bar_y, post_w, bar_h, TEAL)
        if c['post_slot_days'] > 0:
            _add_textbox(slide, bar_left + post_w + Emu(30000), post_bar_y - Emu(10000), Inches(0.4), bar_h + Emu(20000),
                         str(c['post_slot_days']), 'Calibri', 8, GRAY)

    # Footnote
    biggest = max(sorted_cats, key=lambda c: c['change'])
    foot_y = chart_top + group_h * n_cats + Inches(0.15)
    _add_textbox(slide, Inches(0.8), foot_y, Inches(8), Inches(0.3),
                 f"{biggest['category']} coverage claimed {biggest['post_slot_days']} slot-days post-announcement (vs. {biggest['pre_slot_days']} pre), displacing other source types.",
                 'Calibri', 11, GRAY, italic=True)

    _add_logo(slide, logo_path)


def generate_pptx(client_name, event_desc, pivot_date, date_range_str,
                  url_data, cat_summary, n_pre, n_post, summary_points,
                  logo_path, logo_white_path):
    """Build the full 6-slide PPTX deck."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    build_slide_1(prs, client_name, event_desc, date_range_str, logo_white_path)
    build_slide_2(prs, summary_points, logo_path)
    build_slide_3(prs, url_data, n_pre, n_post, client_name, event_desc, pivot_date, logo_path)
    build_slide_4_v2(prs, url_data, n_pre, n_post, logo_path)
    build_slide_5(prs, url_data, n_post, logo_path)
    build_slide_6(prs, cat_summary, n_pre, n_post, logo_path)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────
# AUTHENTICATION
# ─────────────────────────────────────────────

ALLOWED_DOMAIN = "fiveblocks.com"


def _write_google_credentials():
    """Write Google OAuth credentials JSON from Streamlit secrets to a temp file."""
    import tempfile, json
    try:
        client_id = st.secrets["GOOGLE_CLIENT_ID"]
        client_secret = st.secrets["GOOGLE_CLIENT_SECRET"]
    except (KeyError, FileNotFoundError):
        return None

    redirect_uri = _get_redirect_uri()
    creds = {
        "web": {
            "client_id": client_id,
            "client_secret": client_secret,
            "auth_uri": "https://accounts.google.com/o/oauth2/auth",
            "token_uri": "https://oauth2.googleapis.com/token",
            "redirect_uris": [redirect_uri],
        }
    }
    tmp = tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False)
    json.dump(creds, tmp)
    tmp.close()
    return tmp.name


def _get_redirect_uri():
    """Determine redirect URI based on environment."""
    try:
        return st.secrets["REDIRECT_URI"]
    except (KeyError, FileNotFoundError):
        return "http://localhost:8501"


def _get_cookie_key():
    """Get cookie encryption key from secrets, or generate a default for dev."""
    try:
        return st.secrets["COOKIE_KEY"]
    except (KeyError, FileNotFoundError):
        return "serp_evolution_analyzer_dev_key_change_me"


def check_auth():
    """
    Returns True if the user is authenticated with an @fiveblocks.com account.
    Returns False if auth is not configured (local dev — app runs without gate).
    Stops execution (shows login or access denied) if auth is configured but user isn't in.
    """
    # Check if Google OAuth secrets exist
    try:
        _ = st.secrets["GOOGLE_CLIENT_ID"]
        _ = st.secrets["GOOGLE_CLIENT_SECRET"]
    except (KeyError, FileNotFoundError):
        # No OAuth configured — skip auth (local dev mode)
        return True

    from streamlit_google_auth import Authenticate

    creds_path = _write_google_credentials()
    if not creds_path:
        return True

    redirect_uri = _get_redirect_uri()

    auth = Authenticate(
        secret_credentials_path=creds_path,
        redirect_uri=redirect_uri,
        cookie_name="serp_evolution_auth",
        cookie_key=_get_cookie_key(),
    )

    auth.check_authentification()

    if not st.session_state.get('connected', False):
        # Build login URL manually — the google_auth_oauthlib library adds PKCE
        # (code_challenge) which causes 403 errors with Google Workspace Internal apps
        from urllib.parse import urlencode
        params = {
            'response_type': 'code',
            'client_id': st.secrets["GOOGLE_CLIENT_ID"],
            'redirect_uri': redirect_uri,
            'scope': 'openid https://www.googleapis.com/auth/userinfo.profile https://www.googleapis.com/auth/userinfo.email',
            'access_type': 'offline',
            'include_granted_scopes': 'true',
            'hd': ALLOWED_DOMAIN,
        }
        authorization_url = f"https://accounts.google.com/o/oauth2/auth?{urlencode(params)}"

        st.link_button("🔐 Sign in with Google", authorization_url, use_container_width=True)

        st.markdown("---")
        st.markdown(
            "<div style='text-align: center; color: #9CA3AF; margin-top: 2rem;'>"
            "Sign in with your <strong>@fiveblocks.com</strong> Google account to continue."
            "</div>",
            unsafe_allow_html=True
        )
        st.stop()

    # Logged in — check domain
    user_info = st.session_state.get('user_info', {})
    email = user_info.get('email', '')
    if not email.endswith(f'@{ALLOWED_DOMAIN}'):
        st.error(f"⛔ Access denied. This tool is restricted to @{ALLOWED_DOMAIN} accounts.")
        st.caption(f"Signed in as: {email}")
        st.session_state['connected'] = False
        st.session_state['user_info'] = {}
        st.session_state['oauth_id'] = None
        st.stop()

    # Clean up temp file
    try:
        os.unlink(creds_path)
    except Exception:
        pass

    return True


# ─────────────────────────────────────────────
# STREAMLIT APP
# ─────────────────────────────────────────────

def main():
    st.set_page_config(page_title="SERP Evolution Analyzer", page_icon="📊", layout="wide")

    st.markdown("""
    <style>
        .block-container { max-width: 900px; padding-top: 2rem; }
        h1 { color: #1A365D; }
        .stDownloadButton > button { width: 100%; }
    </style>
    """, unsafe_allow_html=True)

    st.title("SERP Evolution Analyzer")
    st.caption("Upload an IMPACT CSV → Get a branded PPTX deck + Excel analysis files")

    # ── Authentication gate ──
    check_auth()

    # Show user info + logout in sidebar if authenticated
    user_info = st.session_state.get('user_info', {})
    is_oauth_active = st.session_state.get('connected', False)

    # Initialize session state
    for key in ['analysis_done', 'deliverables_done', 'url_data', 'n_pre', 'n_post',
                'pre_df', 'post_df', 'cat_summary', 'domain_detail',
                'pptx_bytes', 'xlsx1_bytes', 'xlsx2_bytes']:
        if key not in st.session_state:
            st.session_state[key] = None

    # Sidebar
    with st.sidebar:
        logo_path, logo_white_path = _get_logo_paths()
        if logo_path:
            st.image(logo_path, width=180)
        else:
            st.markdown("**Five Blocks**")

        # Show logged-in user + logout
        if is_oauth_active and user_info:
            st.markdown(f"👤 **{user_info.get('name', '')}**")
            st.caption(user_info.get('email', ''))
            if st.button("Sign out", use_container_width=True):
                st.session_state['connected'] = False
                st.session_state['user_info'] = {}
                st.session_state['oauth_id'] = None
                st.rerun()

        st.markdown("---")

        # API key from Streamlit secrets (invisible to users)
        api_key = None
        try:
            api_key = st.secrets.get("ANTHROPIC_API_KEY", None)
        except Exception:
            pass

        if api_key:
            st.caption("✅ Claude API connected — AI-powered summaries and domain classification enabled.")
        else:
            st.caption("⚡ Running in fallback mode — no API key configured. Domain categorization uses built-in mapping; exec summary uses data-driven templates.")
            st.caption("To enable AI features, add `ANTHROPIC_API_KEY` to Streamlit secrets.")

        st.markdown("---")
        st.markdown("**Five Blocks Internal Tool**")
        st.caption("Upload IMPACT CSV data and configure the analysis parameters.")

    # Main form
    uploaded_file = st.file_uploader("Upload IMPACT CSV", type=['csv'])

    col1, col2 = st.columns(2)
    with col1:
        client_name = st.text_input("Client / Entity name", placeholder="e.g. Connor Teskey")
        pivot_date = st.date_input("Event date (pivot)", value=date(2026, 2, 4))
    with col2:
        event_desc = st.text_input("Event description", placeholder="e.g. CEO succession announcement")
        owned_domains = st.text_input("Owned domains (comma-separated, optional)",
                                      placeholder="e.g. brookfield.com, privatewealth.brookfield.com")

    if not (uploaded_file and client_name and event_desc):
        st.info("Upload an IMPACT CSV and fill in the fields above to get started.")
        return

    file_bytes = uploaded_file.getvalue()

    # ── STEP 1: Run Analysis ──
    if st.button("🔍 Run Analysis", type="primary", use_container_width=True):
        with st.spinner("Processing IMPACT data..."):
            url_data, n_pre, n_post, pre_df, post_df = process_impact_csv(file_bytes, pivot_date)

        with st.spinner("Categorizing domains..."):
            url_data = categorize_domains(url_data, owned_domains, api_key if api_key else None)

        with st.spinner("Computing source type analysis..."):
            cat_summary, domain_detail = compute_source_types(url_data, pre_df, post_df, n_pre, n_post)

        # Store in session state
        st.session_state.url_data = url_data
        st.session_state.n_pre = n_pre
        st.session_state.n_post = n_post
        st.session_state.pre_df = pre_df
        st.session_state.post_df = post_df
        st.session_state.cat_summary = cat_summary
        st.session_state.domain_detail = domain_detail
        st.session_state.analysis_done = True
        st.session_state.deliverables_done = None
        st.rerun()

    # ── Show analysis results if available ──
    if not st.session_state.analysis_done:
        return

    url_data = st.session_state.url_data
    n_pre = st.session_state.n_pre
    n_post = st.session_state.n_post
    pre_df = st.session_state.pre_df
    post_df = st.session_state.post_df
    cat_summary = st.session_state.cat_summary
    domain_detail = st.session_state.domain_detail

    st.success(f"Found {len(url_data)} unique URLs across {n_pre} pre-event and {n_post} post-event days.")

    if abs(n_pre - n_post) > 3:
        st.warning(f"⚠️ Unequal periods: {n_pre} pre-event days vs {n_post} post-event days. Results are still valid but comparisons are not perfectly symmetric.")

    # Domain mapping editor
    with st.expander("📋 Domain → Category Mapping (review & edit)", expanded=False):
        mapping_df = pd.DataFrame([
            {'Domain': d['domain'], 'Category': d['category']}
            for d in url_data.values()
        ]).drop_duplicates().sort_values('Category')

        edited_df = st.data_editor(
            mapping_df,
            column_config={
                "Category": st.column_config.SelectboxColumn(
                    "Category", options=VALID_CATEGORIES, required=True
                )
            },
            hide_index=True, use_container_width=True, key="domain_editor"
        )

        if edited_df is not None:
            edit_map = dict(zip(edited_df['Domain'], edited_df['Category']))
            changed = False
            for url, info in url_data.items():
                if info['domain'] in edit_map and info['category'] != edit_map[info['domain']]:
                    info['category'] = edit_map[info['domain']]
                    changed = True
            if changed:
                cat_summary, domain_detail = compute_source_types(url_data, pre_df, post_df, n_pre, n_post)
                st.session_state.cat_summary = cat_summary
                st.session_state.domain_detail = domain_detail

    # Summary metrics
    n_new = sum(1 for d in url_data.values() if d['status'] == 'New')
    n_persistent = sum(1 for d in url_data.values() if d['status'] == 'Persistent')
    n_dropped = sum(1 for d in url_data.values() if d['status'] == 'Dropped')

    col_a, col_b, col_c, col_d = st.columns(4)
    col_a.metric("Total URLs", len(url_data))
    col_b.metric("New", n_new)
    col_c.metric("Persistent", n_persistent)
    col_d.metric("Dropped", n_dropped)

    # Source type preview
    with st.expander("📊 Source Type Summary", expanded=False):
        st.dataframe(
            pd.DataFrame(cat_summary).rename(columns={
                'category': 'Category', 'pre_slot_days': f'Pre ({n_pre}d)',
                'post_slot_days': f'Post ({n_post}d)', 'change': 'Change'
            }),
            hide_index=True, use_container_width=True
        )

    # ── STEP 2: Generate Deliverables ──
    if st.button("📦 Generate Deliverables", type="primary", use_container_width=True):
        with st.spinner("Generating executive summary..."):
            if api_key:
                try:
                    summary_points = generate_executive_summary(
                        url_data, cat_summary, client_name, event_desc, n_pre, n_post, api_key
                    )
                except Exception as e:
                    st.warning(f"Claude API summary failed ({e}). Using fallback.")
                    summary_points = generate_fallback_summary(
                        url_data, cat_summary, client_name, event_desc, n_pre, n_post
                    )
            else:
                summary_points = generate_fallback_summary(
                    url_data, cat_summary, client_name, event_desc, n_pre, n_post
                )

        first_date = pre_df['parsed_date'].min() if len(pre_df) > 0 else post_df['parsed_date'].min()
        last_date = post_df['parsed_date'].max() if len(post_df) > 0 else pre_df['parsed_date'].max()
        date_range_str = (
            f"Pre- vs. post-{event_desc.lower()}  |  {_format_date(pivot_date, '%b %-d, %Y')}\n"
            f"{_format_date(first_date, '%b %-d, %Y')} – {_format_date(last_date, '%b %-d, %Y')}  ({n_pre} days each side)"
        )

        logo_path, logo_white_path = _get_logo_paths()

        with st.spinner("Building PowerPoint deck..."):
            pptx_bytes = generate_pptx(
                client_name, event_desc, pivot_date, date_range_str,
                url_data, cat_summary, n_pre, n_post, summary_points,
                logo_path, logo_white_path
            )

        with st.spinner("Building Excel files..."):
            xlsx1_bytes = generate_url_analysis_excel(url_data, n_pre, n_post, uploaded_file.name, pivot_date)
            xlsx2_bytes = generate_source_type_excel(cat_summary, domain_detail, n_pre, n_post)

        st.session_state.pptx_bytes = pptx_bytes
        st.session_state.xlsx1_bytes = xlsx1_bytes
        st.session_state.xlsx2_bytes = xlsx2_bytes
        st.session_state.deliverables_done = True
        st.rerun()

    # ── Show downloads if available ──
    if st.session_state.deliverables_done:
        st.markdown("---")
        st.subheader("📥 Downloads")
        safe_name = re.sub(r'[^a-zA-Z0-9_]', '_', client_name)
        dl1, dl2, dl3 = st.columns(3)
        with dl1:
            st.download_button("⬇️ PPTX Deck",
                               st.session_state.pptx_bytes,
                               f"{safe_name}_URL_Evolution.pptx",
                               "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        with dl2:
            st.download_button("⬇️ URL Analysis (Excel)",
                               st.session_state.xlsx1_bytes,
                               f"{safe_name}_Pre_Post_URL_Analysis.xlsx",
                               "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        with dl3:
            st.download_button("⬇️ Source Type Analysis (Excel)",
                               st.session_state.xlsx2_bytes,
                               f"{safe_name}_Source_Type_Analysis.xlsx",
                               "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        st.success("All deliverables generated successfully!")


if __name__ == '__main__':
    main()
